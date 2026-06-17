#!/usr/bin/env python3
"""
Chiron PPT — Full beautification (slides 1-10).
All slides: title locked at y=0.18, content vertically centered, minimal blue-white palette.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

PPTX  = "/Users/rest/Documents/agent-safe/outputs/chiron-report/Chiron-课程项目汇报.pptx"
DARK  = RGBColor(0x0F,0x17,0x2A); BODY = RGBColor(0x33,0x40,0x55)
GRAY  = RGBColor(0x94,0xA3,0xB8); BLUE = RGBColor(0x25,0x63,0xEB)
WHITE = RGBColor(0xFF,0xFF,0xFF); CARD = RGBColor(0xF8,0xFA,0xFC)
BORDER= RGBColor(0xE2,0xE8,0xF0)

TITLE_Y = Inches(0.18)
AVAIL_TOP = 0.70; AVAIL_BOT = 6.90

prs = Presentation(PPTX)

# ─── Helpers ───
def clr(s):
    for sp in list(s.shapes):
        try: sp._element.getparent().remove(sp._element)
        except: pass

def rc(s,x,y,w,h,fill=WHITE,border=BORDER,rounded=True):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,x,y,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = Pt(0.75)
    else: sh.line.fill.background()
    return sh

def tx(s,x,y,w,h,text,sz=Pt(10),c=DARK,b=False,a=PP_ALIGN.LEFT):
    t = s.shapes.add_textbox(x,y,w,h); tf = t.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = sz; p.font.color.rgb = c
    p.font.bold = b; p.font.name = "Arial"; p.alignment = a
    return t

def mtx(s,x,y,w,h,lines,sz=Pt(9),c=BODY,bf=False,a=PP_ALIGN.LEFT):
    t = s.shapes.add_textbox(x,y,w,h); tf = t.text_frame; tf.word_wrap = True
    for i,ln in enumerate(lines):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text = ln; p.font.size = sz; p.font.color.rgb = c
        p.font.name = "Arial"; p.alignment = a
        if bf and i==0: p.font.bold = True
    return t

def bar(s,x,y,w,h):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,x,y,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb = BLUE; sh.line.fill.background()

def arw_r(s,x,y,w):
    sh = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,x,y,w,Inches(0.14))
    sh.fill.solid(); sh.fill.fore_color.rgb = GRAY; sh.line.fill.background()

def arw_d(s,x,y,h):
    sh = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,x,y,Inches(0.18),h)
    sh.fill.solid(); sh.fill.fore_color.rgb = GRAY; sh.line.fill.background()

def pill(s,x,y,w,h,text,fill=BLUE,tc=WHITE):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,x,y,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill; sh.line.fill.background()
    p = sh.text_frame.paragraphs[0]; p.text = text; p.font.size = Pt(8)
    p.font.color.rgb = tc; p.font.bold = True; p.font.name = "Arial"; p.alignment = PP_ALIGN.CENTER

def title(s,text):
    tx(s,Inches(0.50),TITLE_Y,Inches(12.0),Inches(0.42),text,Pt(23),DARK,True)

def coff(h):
    return (AVAIL_BOT - AVAIL_TOP - h) / 2


# ═══════════════════════════════════════════
# SLIDE 1 — COVER
# ═══════════════════════════════════════════
s = prs.slides[0]; clr(s)

# Full-slide dark background
bg = rc(s,Inches(0),Inches(0),Inches(13.33),Inches(7.50),DARK,None,False)

# Center content block
y0 = Inches(1.0)
tx(s,Inches(1.0),y0,Inches(11.0),Inches(0.80),"Chiron",Pt(48),WHITE,True,PP_ALIGN.CENTER)
tx(s,Inches(1.0),y0+Inches(0.85),Inches(11.0),Inches(0.40),
    "Agent 交易安全中间件  |  Intent-Execution Semantic Consistency Verification",
    Pt(14),GRAY,False,PP_ALIGN.CENTER)

# Accent line
rc(s,Inches(5.0),y0+Inches(1.40),Inches(3.5),Inches(0.04),BLUE,None,False)

# Quote
tx(s,Inches(1.5),y0+Inches(1.70),Inches(10.0),Inches(0.90),
    "「Agent 说：Swap 100 USDC for ETH on Uniswap V3\n"
    "  Chiron 检查：实际 calldata 真的是在 Swap 100 USDC 吗？\n"
    "  结果：一致 → 放行  |  不一致 → 阻止 + 告警」",
    Pt(13),RGBColor(0xCB,0xD5,0xE1),False,PP_ALIGN.CENTER)

# Tags
tags = ["AI × 区块链安全","语义一致性验证","前置阻断架构"]
for i,tg in enumerate(tags):
    pill(s,Inches(3.2+i*2.3),y0+Inches(2.85),Inches(2.0),Inches(0.35),tg)

# Footer
tx(s,Inches(1.0),Inches(6.60),Inches(11.0),Inches(0.30),
    "课程项目汇报  |  2026.06",Pt(11),GRAY,False,PP_ALIGN.CENTER)


# ═══════════════════════════════════════════
# SLIDE 2
# ═══════════════════════════════════════════
s = prs.slides[1]; clr(s)
title(s,"Agent 交易的全过程——「语义鸿沟」在哪里")
off = Inches(AVAIL_TOP + coff(5.05))

steps = [
    ("① 预设策略触发","ETH < $1600 → 条件满足\nAgent 决定用 1000 USDC 换 WETH"),
    ("② LLM 生成指令","Uniswap V3 exactInputSingle\n(USDC → WETH, 1000e6)"),
    ("③ 构建 calldata 并签名","0x414bf389... 签名 → 广播至 Mempool"),
]
for i,(ti,d) in enumerate(steps):
    sx = Inches(0.50)+Inches(i*4.20)
    rc(s,sx,off,Inches(3.80),Inches(1.45)); bar(s,sx,off,Inches(3.80),Inches(0.04))
    tx(s,sx+Inches(0.20),off+Inches(0.13),Inches(3.40),Inches(0.28),ti,Pt(13),DARK,True)
    tx(s,sx+Inches(0.20),off+Inches(0.50),Inches(3.40),Inches(0.80),d,Pt(11),BODY)
    if i<2: arw_r(s,sx+Inches(3.80),off+Inches(0.55),Inches(0.40))

wy = off+Inches(1.65)
rc(s,Inches(0.50),wy,Inches(12.30),Inches(0.60),CARD,BORDER)
tx(s,Inches(0.80),wy+Inches(0.10),Inches(11.70),Inches(0.42),
    "⚠ 关键空白：步骤②→③之间无一致性检查。Agent「声称要做的事」和「实际签署的事」可能完全不同。",Pt(12),DARK,True)

cy = off+Inches(2.45)
rc(s,Inches(0.50),cy,Inches(5.95),Inches(2.30)); bar(s,Inches(0.50),cy,Inches(5.95),Inches(0.04))
tx(s,Inches(0.80),cy+Inches(0.15),Inches(5.40),Inches(0.28),"Agent 声明的意图",Pt(14),DARK,True)
mtx(s,Inches(0.80),cy+Inches(0.55),Inches(5.40),Inches(1.60),
    ["动作：Swap          协议：Uniswap V3","输入：1000 USDC → WETH",
     "接收：0xUserWallet   滑点：0.5%","","IntentTemplate → keccak256 → intentHash"],Pt(12),BODY)
tx(s,Inches(6.50),cy+Inches(0.85),Inches(0.60),Inches(0.40),"VS",Pt(18),GRAY,True,PP_ALIGN.CENTER)
rc(s,Inches(7.20),cy,Inches(5.95),Inches(2.30)); bar(s,Inches(7.20),cy,Inches(5.95),Inches(0.04))
tx(s,Inches(7.50),cy+Inches(0.15),Inches(5.40),Inches(0.28),"Prompt 注入后实际签署的 calldata",Pt(14),DARK,True)
mtx(s,Inches(7.50),cy+Inches(0.55),Inches(5.40),Inches(1.60),
    ["目标：0xA0b8...USDC（非 Uniswap Router）","selector：0x095ea7b3 → approve",
     "spender：0xAttacker","amount：0xffff...（无限授权）"],Pt(12),BODY)
tx(s,Inches(0.50),off+Inches(4.95),Inches(12.30),Inches(0.30),
    "Chiron 的解法：在步骤②→③之间插入一致性检查，比对 calldata 与 IntentTemplate，不一致则阻止签名。",Pt(10),GRAY,False,PP_ALIGN.CENTER)


# ═══════════════════════════════════════════
# SLIDE 3
# ═══════════════════════════════════════════
s = prs.slides[2]; clr(s)
title(s,"Prompt 注入——Agent 交易面临的最大威胁")
off = Inches(AVAIL_TOP + coff(4.15))
tx(s,Inches(0.50),off-Inches(0.10),Inches(12.0),Inches(0.25),
    "攻击者将恶意指令嵌入 Agent 接收的外部数据中，AI 无法区分「数据」与「指令」",Pt(11),GRAY)
paths = [
    ("🔗 路径一：链上数据注入","Agent 查询链上数据时，攻击者在合约\n日志中嵌入恶意指令。\n\nPenligent.ai (2026) CCSC 攻击\n影响约 $500M 跨链桥资产 [6]"),
    ("🌐 路径二：API 响应操控","Agent 依赖的链下 API 响应被篡改。\n\nOpenAI EVMbench (2025)：\n主流 LLM 正确识别注入仅 12% [4]"),
    ("🧠 路径三：记忆/上下文污染","Agent 记忆被注入虚假记录。\n\nRyder.id (2026)：\n上下文操纵攻击增长 300%+ [7]"),
]
for i,(ti,d) in enumerate(paths):
    px = Inches(0.50)+Inches(i*4.20)
    rc(s,px,off+Inches(0.20),Inches(3.80),Inches(2.70)); bar(s,px,off+Inches(0.20),Inches(3.80),Inches(0.04))
    tx(s,px+Inches(0.20),off+Inches(0.32),Inches(3.40),Inches(0.30),ti,Pt(13),DARK,True)
    tx(s,px+Inches(0.20),off+Inches(0.75),Inches(3.40),Inches(2.00),d,Pt(10),BODY)

dy = off+Inches(3.10)
rc(s,Inches(0.50),dy,Inches(12.30),Inches(0.95),CARD,BORDER)
mtx(s,Inches(0.80),dy+Inches(0.10),Inches(11.70),Inches(0.78),
    ["Chiron 的防御：在签名前发现 intent(swap) ≠ calldata(approve)，阻止签名。",
     "不依赖 AI 判断安全性——只做确定性的一致性比对。Prompt 注入可以欺骗 LLM，但改不了 calldata 的结构。",
     "R1（目标合约检查）+ R2（selector 检查）即可覆盖以上全部三条路径。"],Pt(10),BODY,True)
tx(s,Inches(0.50),off+Inches(4.20),Inches(12.30),Inches(0.25),
    "参考：[1] arXiv 2601.04583  [4] OpenAI EVMbench  [6] Penligent.ai  [7] Ryder.id",Pt(8),GRAY)


# ═══════════════════════════════════════════
# SLIDE 4
# ═══════════════════════════════════════════
s = prs.slides[3]; clr(s)
title(s,"更多攻击场景——Chiron 全覆盖")
off = Inches(AVAIL_TOP + coff(4.90))

attacks = [
    ("📬 接收地址替换","Agent 发送 WETH 到 0xUser\n攻击者篡改接收参数为 0xAttacker [5]","R5: receiver 不匹配"),
    ("🕵 隐藏 delegatecall","正常合约通过 delegatecall\n跳转到攻击者恶意代码 [1]","R6: 检测 delegatecall"),
    ("💰 金额操纵 10x","amountIn 从 100e6 改为\n1000e6（10 倍）[2]","R4: 超出滑点范围"),
    ("🎣 钓鱼合约（1 字符差）","虚假 Router 地址\n仅最后 1 位不同 [3]","R1: 不在已知协议表"),
]
for i,(ti,d,df) in enumerate(attacks):
    r,c = i//2, i%2
    ax,ay = Inches(0.50)+Inches(c*6.30), off+Inches(r*2.30)
    rc(s,ax,ay,Inches(5.90),Inches(2.00)); bar(s,ax,ay,Inches(5.90),Inches(0.04))
    tx(s,ax+Inches(0.25),ay+Inches(0.15),Inches(3.50),Inches(0.30),ti,Pt(14),DARK,True)
    tx(s,ax+Inches(0.25),ay+Inches(0.55),Inches(3.50),Inches(0.80),d,Pt(11),BODY)
    rc(s,ax+Inches(3.90),ay+Inches(0.30),Inches(1.75),Inches(0.55),CARD,BORDER)
    tx(s,ax+Inches(4.00),ay+Inches(0.40),Inches(1.55),Inches(0.40),df,Pt(10),BLUE,True,PP_ALIGN.CENTER)
tx(s,Inches(0.50),off+Inches(4.80),Inches(12.30),Inches(0.25),
    "参考：[1] arXiv 2601.04583  [2] arXiv 2602.17805  [3] TRM Labs  [5] OWASP",Pt(8),GRAY)


# ═══════════════════════════════════════════
# SLIDE 5
# ═══════════════════════════════════════════
s = prs.slides[4]; clr(s)
title(s,"现有方案为什么挡不住这些攻击？")
off = Inches(AVAIL_TOP + coff(5.05))
tx(s,Inches(0.50),off-Inches(0.10),Inches(12.0),Inches(0.25),
    "所有方案都在回答「交易安全吗？」——一个不可判定问题",Pt(11),GRAY)

sols = [
    ("Transaction Simulation  (Blowfish / Tenderly)","模拟交易 → Agent 自动执行，无人看屏幕 [8]"),
    ("On-chain Monitoring  (Forta / Dune)","上链后告警 → approve 已确认，资金已转走 [2]"),
    ("Policy Engine  (Zodiac / ERC-6900)","白名单只检查地址，不检查语义一致性 [5]"),
    ("Multi-sig  (Safe / Gnosis)","每笔交易等人签字，违背 Agent 自动化目的 [4]"),
]
for i,(nm,rs) in enumerate(sols):
    sy = off+Inches(i*1.10)
    rc(s,Inches(0.50),sy,Inches(12.30),Inches(0.90)); bar(s,Inches(0.50),sy,Inches(0.05),Inches(0.90))
    tx(s,Inches(0.80),sy+Inches(0.12),Inches(4.00),Inches(0.60),nm,Pt(12),DARK,True)
    tx(s,Inches(5.00),sy+Inches(0.12),Inches(5.00),Inches(0.60),rs,Pt(11),BODY)
    tx(s,Inches(10.50),sy+Inches(0.20),Inches(1.50),Inches(0.40),"✕ 无法阻止",Pt(11),GRAY,True,PP_ALIGN.CENTER)

by = off+Inches(4.60)
rc(s,Inches(0.50),by,Inches(12.30),Inches(0.65),CARD,BORDER)
mtx(s,Inches(0.80),by+Inches(0.10),Inches(11.70),Inches(0.48),
    ["Chiron 换了一个问题：「交易和 Agent 声称的意图一致吗？」→ 确定性可验证，10ms，0 Gas",
     "不需要知道「什么是对的」，只需要知道「是否一致」——将不可判定问题转化为可判定问题。"],Pt(11),BODY,True)


# ═══════════════════════════════════════════
# SLIDE 6
# ═══════════════════════════════════════════
s = prs.slides[5]; clr(s)
title(s,"Chiron 的回答：从「Is it safe?」到「Is it consistent?」")
off = Inches(AVAIL_TOP + coff(5.05))

cy = off
rc(s,Inches(0.50),cy,Inches(5.80),Inches(1.85)); bar(s,Inches(0.50),cy,Inches(5.80),Inches(0.04))
tx(s,Inches(0.80),cy+Inches(0.13),Inches(5.20),Inches(0.28),"旧问题：Is it safe?",Pt(15),DARK,True)
mtx(s,Inches(0.80),cy+Inches(0.50),Inches(5.20),Inches(1.20),
    ["要回答「交易安全吗」必须先知道用户真实意图。","用户意图在心智中，无法被形式化证明。",
     "","→ 所有方案困在不可判定问题 [1]"],Pt(12),BODY)
tx(s,Inches(6.40),cy+Inches(0.68),Inches(0.50),Inches(0.40),"⟶",Pt(28),GRAY,True,PP_ALIGN.CENTER)
rc(s,Inches(7.00),cy,Inches(5.80),Inches(1.85)); bar(s,Inches(7.00),cy,Inches(5.80),Inches(0.04))
tx(s,Inches(7.30),cy+Inches(0.13),Inches(5.20),Inches(0.28),"新问题：Is it consistent?",Pt(15),DARK,True)
mtx(s,Inches(7.30),cy+Inches(0.50),Inches(5.20),Inches(1.20),
    ["只需比较两个结构化声明：","IntentTemplate vs DecodedTx","","→ 确定性可验证，10ms，0 Gas"],Pt(12),BODY)

fy = off+Inches(2.10)
tx(s,Inches(0.50),fy,Inches(12.0),Inches(0.32),"验证流程示例（Prompt 注入）",Pt(14),DARK,True)

flow = [
    ("步骤 1","Agent 声明意图","chiron.intent(\"swap\",...)\n→ IntentTemplate { SWAP }"),
    ("步骤 2","Chiron 解码交易","TxDecoder.decode(...)\n→ DecodedTx { approve }"),
    ("步骤 3","L1 规则检查","R1: target ≠ Router → FAIL\nR2: selector ≠ swap → FAIL\n→ 阻止签名"),
]
for i,(st,ti,de) in enumerate(flow):
    fx,fyy = Inches(0.50)+Inches(i*4.20), fy+Inches(0.40)
    rc(s,fx,fyy,Inches(3.80),Inches(1.85)); bar(s,fx,fyy,Inches(3.80),Inches(0.04))
    pill(s,fx+Inches(0.15),fyy+Inches(0.13),Inches(0.90),Inches(0.26),st)
    tx(s,fx+Inches(1.20),fyy+Inches(0.13),Inches(2.40),Inches(0.26),ti,Pt(13),DARK,True)
    tx(s,fx+Inches(0.20),fyy+Inches(0.53),Inches(3.40),Inches(1.20),de,Pt(11),BODY)
    if i<2: arw_r(s,fx+Inches(3.80),fyy+Inches(0.75),Inches(0.40))

sy = fy+Inches(2.50)
rc(s,Inches(0.50),sy,Inches(12.30),Inches(0.50),CARD,BORDER)
tx(s,Inches(0.80),sy+Inches(0.09),Inches(11.70),Inches(0.32),
    "兼容标准：ERC-8004 (Trustless Agents)  |  ERC-7683 (Cross Chain Intents)  |  ERC-8211 (Smart Batching)  |  ERC-4337 (Account Abstraction)",
    Pt(10),GRAY,False,PP_ALIGN.CENTER)


# ═══════════════════════════════════════════
# SLIDE 7 — Architecture
# ═══════════════════════════════════════════
s = prs.slides[6]; clr(s)
tx(s,Inches(0.40),TITLE_Y,Inches(12.5),Inches(0.42),"Chiron 验证架构——完整数据流",Pt(23),DARK,True)

ROW1_Y=Inches(AVAIL_TOP-0.145); ROW1_H=Inches(0.82)
ROW2_Y=ROW1_Y+ROW1_H+Inches(0.20); ROW2_H=Inches(1.72)
ROW3_Y=ROW2_Y+ROW2_H+Inches(0.20); ROW3_H=Inches(1.30)
ROW4_Y=ROW3_Y+ROW3_H+Inches(0.20); ROW4_H=Inches(1.25)
MET_Y=ROW4_Y+ROW4_H+Inches(0.20); MET_H=Inches(0.60)

# Row 1
rc(s,Inches(0.40),ROW1_Y,Inches(12.50),ROW1_H,CARD,BORDER)
tx(s,Inches(0.55),ROW1_Y+Inches(0.03),Inches(2.5),Inches(0.20),"① Agent 输入",Pt(8),GRAY,True)
rc(s,Inches(1.80),ROW1_Y+Inches(0.12),Inches(4.40),Inches(0.58)); bar(s,Inches(1.80),ROW1_Y+Inches(0.12),Inches(4.40),Inches(0.04))
mtx(s,Inches(1.95),ROW1_Y+Inches(0.15),Inches(4.10),Inches(0.52),
    ["IntentTemplate（Agent 声明的意图）","action:swap  tokenIn:USDC  amount:100e6  → keccak256 → intentHash"],Pt(9),DARK,True)
arw_r(s,Inches(6.25),ROW1_Y+Inches(0.32),Inches(0.80))
rc(s,Inches(7.10),ROW1_Y+Inches(0.12),Inches(4.40),Inches(0.58)); bar(s,Inches(7.10),ROW1_Y+Inches(0.12),Inches(4.40),Inches(0.04))
mtx(s,Inches(7.25),ROW1_Y+Inches(0.15),Inches(4.10),Inches(0.52),
    ["TxCandidate（待签名交易）","target:0xE592...Router  selector:0x414bf389  → TxDecoder 解码"],Pt(9),DARK,True)
arw_d(s,Inches(6.60),ROW1_Y+ROW1_H-Inches(0.02),Inches(0.18))

# Row 2
rc(s,Inches(0.40),ROW2_Y,Inches(12.50),ROW2_H,CARD,BORDER)
tx(s,Inches(0.55),ROW2_Y+Inches(0.03),Inches(7.0),Inches(0.22),"② L1: 确定性本地校验（主路径·0 Gas·~10ms·覆盖 95% 交易）",Pt(8),GRAY,True)
rc(s,Inches(1.20),ROW2_Y+Inches(0.24),Inches(7.20),Inches(1.28)); bar(s,Inches(1.20),ROW2_Y+Inches(0.24),Inches(7.20),Inches(0.04))
mtx(s,Inches(1.40),ROW2_Y+Inches(0.32),Inches(3.30),Inches(1.10),
    ["R1 目标合约 ∈ 已知协议地址表","R2 函数 selector 匹配 Action 类型","R3 TokenIn / TokenOut 地址一致"],Pt(9.5),DARK)
mtx(s,Inches(4.80),ROW2_Y+Inches(0.32),Inches(3.40),Inches(1.10),
    ["R4 金额在 ±0.5% 滑点范围内","R5 无意外参数（receiver/approve）","R6 无嵌套 delegatecall / CREATE2"],Pt(9.5),DARK)
pill(s,Inches(3.00),ROW2_Y+Inches(0.26),Inches(1.20),Inches(0.18),"~10ms · 0 Gas")
pill(s,Inches(5.50),ROW2_Y+Inches(1.26),Inches(1.00),Inches(0.18),"95% 覆盖")
arw_r(s,Inches(8.45),ROW2_Y+Inches(0.66),Inches(0.65))
tx(s,Inches(8.55),ROW2_Y+Inches(0.40),Inches(0.50),Inches(0.22),"PASS",Pt(9),BLUE,True,PP_ALIGN.CENTER)
rc(s,Inches(9.20),ROW2_Y+Inches(0.24),Inches(2.60),Inches(0.60),CARD,BORDER)
mtx(s,Inches(9.35),ROW2_Y+Inches(0.30),Inches(2.30),Inches(0.48),["验证通过","→ Agent 签名 → 广播上链"],Pt(10),BODY,True,PP_ALIGN.CENTER)
tx(s,Inches(3.80),ROW2_Y+Inches(1.45),Inches(2.00),Inches(0.22),"FAIL / UNCERTAIN (5%)",Pt(8),GRAY,True,PP_ALIGN.CENTER)
arw_d(s,Inches(4.65),ROW2_Y+ROW2_H-Inches(0.02),Inches(0.18))
rc(s,Inches(9.20),ROW2_Y+Inches(1.02),Inches(2.60),Inches(0.48),CARD,BORDER)
mtx(s,Inches(9.35),ROW2_Y+Inches(1.08),Inches(2.30),Inches(0.38),["直接阻止","→ 触发 CircuitBreaker"],Pt(9),BODY,True,PP_ALIGN.CENTER)

# Row 3
rc(s,Inches(0.40),ROW3_Y,Inches(12.50),ROW3_H,CARD,BORDER)
tx(s,Inches(0.55),ROW3_Y+Inches(0.03),Inches(7.0),Inches(0.22),"③ L2: 多模型交叉验证（仅 L1 无法判定时触发·非阻塞异步路径）",Pt(8),GRAY,True)
rc(s,Inches(1.20),ROW3_Y+Inches(0.24),Inches(7.20),Inches(0.88)); bar(s,Inches(1.20),ROW3_Y+Inches(0.24),Inches(7.20),Inches(0.04))
for i,md in enumerate(["验证器 A","验证器 B","验证器 C"]):
    mx = Inches(1.50)+Inches(i*2.35)
    rc(s,mx,ROW3_Y+Inches(0.32),Inches(2.05),Inches(0.64))
    tx(s,mx+Inches(0.05),ROW3_Y+Inches(0.34),Inches(1.95),Inches(0.28),md,Pt(12),DARK,True,PP_ALIGN.CENTER)
    tx(s,mx+Inches(0.05),ROW3_Y+Inches(0.66),Inches(1.95),Inches(0.22),"独立语义判定",Pt(8),GRAY,False,PP_ALIGN.CENTER)
tx(s,Inches(2.00),ROW3_Y+Inches(1.02),Inches(5.00),Inches(0.18),"多个独立模型交叉验证·多数一致即通过    |    不一致→阻止",Pt(8),GRAY,False,PP_ALIGN.CENTER)
arw_r(s,Inches(8.45),ROW3_Y+Inches(0.46),Inches(0.65))
tx(s,Inches(8.55),ROW3_Y+Inches(0.22),Inches(0.50),Inches(0.22),"PASS",Pt(9),BLUE,True,PP_ALIGN.CENTER)
rc(s,Inches(9.20),ROW3_Y+Inches(0.24),Inches(2.60),Inches(0.42),CARD,BORDER)
tx(s,Inches(9.35),ROW3_Y+Inches(0.30),Inches(2.30),Inches(0.28),"→ Agent 签名广播",Pt(10),BODY,True,PP_ALIGN.CENTER)
tx(s,Inches(9.80),ROW3_Y+Inches(0.70),Inches(1.40),Inches(0.22),"FAIL",Pt(10),GRAY,True,PP_ALIGN.CENTER)
rc(s,Inches(9.20),ROW3_Y+Inches(0.72),Inches(2.60),Inches(0.48),CARD,BORDER)
mtx(s,Inches(9.35),ROW3_Y+Inches(0.78),Inches(2.30),Inches(0.38),["阻止签名","→ 告警 + 链上记录"],Pt(9),BODY,True,PP_ALIGN.CENTER)

# Row 4
rc(s,Inches(0.40),ROW4_Y,Inches(12.50),ROW4_H,CARD,BORDER)
tx(s,Inches(0.55),ROW4_Y+Inches(0.03),Inches(8.0),Inches(0.22),"④ 链上合约层（EVM 安全网·异步存证·经济质押）",Pt(8),GRAY,True)
for i,(ti,d) in enumerate([("CircuitBreaker","连续 N 次 FAIL\n自动暂停 Agent"),("BondPool","Agent 质押 ERC-20\n质押量 → 交易上限"),
    ("VerificationStore","验证结果哈希上链\n支持查询与挑战"),("IntentRegistry","协议地址 → Action 映射\n社区贡献·可扩展")]):
    cx = Inches(0.60)+Inches(i*3.15)
    rc(s,cx,ROW4_Y+Inches(0.28),Inches(2.90),Inches(0.84)); bar(s,cx,ROW4_Y+Inches(0.28),Inches(2.90),Inches(0.04))
    tx(s,cx+Inches(0.15),ROW4_Y+Inches(0.38),Inches(2.60),Inches(0.24),ti,Pt(10),DARK,True)
    tx(s,cx+Inches(0.15),ROW4_Y+Inches(0.66),Inches(2.60),Inches(0.42),d,Pt(8.5),BODY)

# Metrics
rc(s,Inches(0.40),MET_Y,Inches(12.50),MET_H,DARK,None,False)
for i,(v,l) in enumerate([("~10ms","L1 延迟"),("0 Gas","Gas 开销"),("~95%","交易覆盖率"),("40+","内置协议")]):
    mx = Inches(0.60)+Inches(i*3.15)
    tx(s,mx,MET_Y+Inches(0.04),Inches(2.90),Inches(0.30),v,Pt(16),WHITE,True,PP_ALIGN.CENTER)
    tx(s,mx,MET_Y+Inches(0.36),Inches(2.90),Inches(0.20),l,Pt(8),GRAY,False,PP_ALIGN.CENTER)


# ═══════════════════════════════════════════
# SLIDE 8 — Attack Coverage
# ═══════════════════════════════════════════
s = prs.slides[7]; clr(s)
title(s,"攻防数据：15 种攻击类型全覆盖")
off = Inches(AVAIL_TOP + coff(5.60))
tx(s,Inches(0.50),off-Inches(0.10),Inches(12.0),Inches(0.25),
    "47 测试通过·15/15 攻击类型覆盖·9/9 L1 规则有效阻断",Pt(11),GRAY)

attacks = [
    ("A1","Prompt 注入交易替换","R1+R2","l1"),("A4","隐藏 Approve 调用","R5","l1"),("A5","隐藏 delegatecall","R6","l1"),
    ("A6","金额操纵 10x","R4","l1"),("A7","接收地址替换","R5","l1"),("A8","钓鱼合约（地址差 1 位）","R1","l1"),
    ("A10","私钥泄露后滥用","L1 二次防线","chain"),("A13","验证证明重放","链上 txHash","chain"),("A14","日交易额度耗尽","合约限制","chain"),
]
CW,CH,GX,GY = Inches(3.80), Inches(1.35), Inches(0.40), Inches(0.25)
for idx,(aid,nm,ru,fl) in enumerate(attacks):
    r,c = idx//3, idx%3
    cx,cy = Inches(0.50)+c*(CW+GX), off+r*(CH+GY)
    rc(s,cx,cy,CW,CH); bar(s,cx,cy,CW,Inches(0.04))
    pc = BLUE if fl=="l1" else GRAY
    pill(s,cx+Inches(0.15),cy+Inches(0.18),Inches(0.55),Inches(0.28),aid,fill=pc)
    pill(s,cx+Inches(0.80),cy+Inches(0.18),Inches(1.40),Inches(0.28),"→ "+ru,fill=pc)
    tx(s,cx+Inches(0.15),cy+Inches(0.60),CW-Inches(0.30),Inches(0.55),nm,Pt(13),DARK,True)

by = off+Inches(3*CH+2*GY)
rc(s,Inches(0.50),by,Inches(12.30),Inches(0.65),CARD,BORDER)
tb = s.shapes.add_textbox(Inches(0.80),by+Inches(0.08),Inches(11.70),Inches(0.50))
tf = tb.text_frame; tf.word_wrap = True
p0 = tf.paragraphs[0]
p0.text = "AttackTestRunner 自动化测试框架：构造攻击向量 → L1 一致性校验 → 阻断验证。L1 规则覆盖 9 种注入/操纵类攻击（9/9 有效阻断）。"
p0.font.size = Pt(10); p0.font.color.rgb = BODY; p0.font.name = "Arial"
p1 = tf.add_paragraph()
p1.text = "A15（恶意 Intent）为设计已知边界。链上合约层（A10/A13/A14）提供经济质押+存证防重放+额度限制三道防线。"
p1.font.size = Pt(10); p1.font.color.rgb = GRAY; p1.font.name = "Arial"


# ═══════════════════════════════════════════
# SLIDE 9 — References
# ═══════════════════════════════════════════
s = prs.slides[8]; clr(s)
title(s,"参考文献与论据来源")

refs = [
    "[1] arXiv 2601.04583 (2025) —「Autonomous Agents on Blockchains」— 形式化描述意图-执行鸿沟 (§5)",
    "[2] arXiv 2602.17805 (2026) —「Liquidity Exhaustion Attacks」— 分析 350 万笔跨链意图，deBridge 攻击成功率 80.5%",
    "[3] TRM Labs / Chainalysis (2026) —「2026 Crypto Crime Report」— AI 增强攻击损失同比增长 40%",
    "[4] OpenAI (2025) — EVMbench — 主流 LLM 正确识别注入仅 12%",
    "[5] OWASP (2026) —「Smart Contract Top 10」— 策略绕过列 #4",
    "[6] Penligent.ai (2026) —「The 2026 Sovereign Ledger」— CCSC 攻击影响约 $500M",
    "[7] Ryder.id (2026) —「AI-Powered Crypto Attacks in 2026」— 上下文操纵攻击增长 300%+",
    "[8] ERC-8004 (2025) — Trustless Agents (MetaMask/Coinbase/EF) — §7：当前方案「代码即意图假设」",
]

off = Inches(AVAIL_TOP + coff(len(refs)*0.55 + 0.40))
for i,r in enumerate(refs):
    ry = off + Inches(i*0.55)
    rc(s,Inches(0.50),ry,Inches(12.30),Inches(0.45)); bar(s,Inches(0.50),ry,Inches(0.05),Inches(0.45))
    tx(s,Inches(0.75),ry+Inches(0.08),Inches(11.80),Inches(0.30),r,Pt(10),BODY)


# ═══════════════════════════════════════════
# SLIDE 10 — Thank You
# ═══════════════════════════════════════════
s = prs.slides[9]; clr(s)

# Dark background
rc(s,Inches(0),Inches(0),Inches(13.33),Inches(7.50),DARK,None,False)

# Thank You
tx(s,Inches(1.0),Inches(1.20),Inches(11.0),Inches(0.80),"Thank You",Pt(48),WHITE,True,PP_ALIGN.CENTER)
tx(s,Inches(1.0),Inches(1.95),Inches(11.0),Inches(0.35),
    "Chiron — Agent Transaction Security Middleware",Pt(14),GRAY,False,PP_ALIGN.CENTER)

# Divider
rc(s,Inches(5.0),Inches(2.50),Inches(3.5),Inches(0.04),BLUE,None,False)

# Quote
tx(s,Inches(2.0),Inches(2.80),Inches(9.0),Inches(0.60),
    "「一致性是可验证的，正确性是不可验证的。\n  Chiron 保证一致性，不保证正确性。」",
    Pt(14),RGBColor(0xCB,0xD5,0xE1),False,PP_ALIGN.CENTER)

# Philosophy box
rc(s,Inches(2.0),Inches(3.70),Inches(9.0),Inches(1.60),
    RGBColor(0x1E,0x29,0x3B),RGBColor(0x33,0x40,0x55))
tx(s,Inches(2.30),Inches(3.85),Inches(8.40),Inches(0.25),
    "Chiron 的安全模型：一致性 > 正确性",Pt(14),WHITE,True,PP_ALIGN.CENTER)
mtx(s,Inches(2.30),Inches(4.20),Inches(8.40),Inches(0.90),
    ["一致性 = 交易执行内容 = 声明要做的事（可确定性验证）",
     "正确性 = 交易是否安全/理性（不可形式化验证）",
     "",
     "一个可验证的「一致」，比一个无法验证的「安全」更有价值。"],
    Pt(11),GRAY,False,PP_ALIGN.CENTER)

# Footer
tx(s,Inches(1.0),Inches(6.50),Inches(11.0),Inches(0.30),
    "github.com/agent-safe  |  npm: @chiron/sdk",Pt(10),GRAY,False,PP_ALIGN.CENTER)


# ═══════════════════════════════════════════
prs.save(PPTX)
print("✅ All 10 slides beautified with unified style")
